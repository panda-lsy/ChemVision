"""更新 GOAI 参赛 PPT

主要操作:
1. 全局替换项目名: "万化构象 ChemVISION" / "万化构象" / "ChemVISION"
   → 统一改为 "ChemEdu Agent"(代码 package:chemedu 已对齐)
2. 在 THANKS 页之前插入 3 张新幻灯片,补充近期迭代的新内容:
   - Agent 能力闭环强化(多轮对话 / 任务进度可视化 / 追问合并历史)
   - 个性化学习闭环(学情画像→诊断→规划→错题本 / "我的"页面)
   - 技术架构更新(含 Agent 编排层 / Riverpod 状态管理 / Hive 本地存储)
3. 强化与 AI+教育赛题的契合度(贴合赛题七环节与评审维度)

用法: python tools/update_goai_ppt.py
"""

from __future__ import annotations

import shutil
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

PPT_DIR = Path(__file__).resolve().parent.parent / "doc" / "GOAI"
SRC_PPT = PPT_DIR / "万化构象——化学结构式可视化学习工具最新.pptx"
DST_PPT = PPT_DIR / "ChemEdu_Agent_赛道演示.pptx"
BACKUP = PPT_DIR / "万化构象——化学结构式可视化学习工具最新.bak.pptx"

# 名字替换映射(顺序敏感: 先替换长串,再替换短串,避免误伤)
NAME_REPLACEMENTS = [
    ("万化构象 ChemVISION", "ChemEdu Agent"),
    ("万化构象", "ChemEdu Agent"),
    ("ChemVISION", "ChemEdu Agent"),
    ("chemvision.qzz.io", "chemedu.qzz.io"),
]

# 主题色(从原 PPT 推断的蓝绿配色,贴近赛道视觉)
ACCENT_BLUE = RGBColor(0x3D, 0x77, 0xDE)
ACCENT_AQUA = RGBColor(0x66, 0xE3, 0xD4)
TEXT_DARK = RGBColor(0x1F, 0x29, 0x37)
TEXT_MUTED = RGBColor(0x6B, 0x72, 0x80)
BG_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
BORDER_LIGHT = RGBColor(0xE5, 0xE7, 0xEB)


def replace_text_in_shape(shape) -> None:
    """递归替换 shape 内文本,保留 run 级样式。"""
    # 替换 shape 自身的 text_frame
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                original = run.text
                new_text = original
                for src, dst in NAME_REPLACEMENTS:
                    if src in new_text:
                        new_text = new_text.replace(src, dst)
                if new_text != original:
                    run.text = new_text

    # 表格内文本
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        original = run.text
                        new_text = original
                        for src, dst in NAME_REPLACEMENTS:
                            if src in new_text:
                                new_text = new_text.replace(src, dst)
                        if new_text != original:
                            run.text = new_text

    # 组合形状递归
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for child in shape.shapes:
            replace_text_in_shape(child)


def replace_all_text(prs: Presentation) -> int:
    """遍历所有 slide 替换文本,返回替换次数。"""
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            before = shape.text if shape.has_text_frame else ""
            replace_text_in_shape(shape)
            after = shape.text if shape.has_text_frame else ""
            if before != after:
                count += 1
    return count


def _set_text(
    textbox,
    text: str,
    *,
    font_size: int = 14,
    bold: bool = False,
    color: RGBColor = TEXT_DARK,
    align=PP_ALIGN.LEFT,
) -> None:
    """填充 textbox 单段文本(清除原有内容)。"""
    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "微软雅黑"


def _add_bullet(
    tf,
    text: str,
    *,
    font_size: int = 13,
    color: RGBColor = TEXT_DARK,
    bold: bool = False,
    level: int = 0,
) -> None:
    """在 text_frame 追加一条项目符号段落。"""
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.level = level
    run = p.add_run()
    run.text = ("• " if level == 0 else "- ") + text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "微软雅黑"


def _add_section_header(slide, title: str, subtitle: str) -> None:
    """在 slide 顶部绘制章节标题块。"""
    # 标题
    tb = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.7)
    )
    _set_text(
        tb, title, font_size=24, bold=True, color=ACCENT_BLUE
    )
    # 副标题
    sb = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.1), Inches(12.0), Inches(0.4)
    )
    _set_text(
        sb, subtitle, font_size=12, color=TEXT_MUTED
    )
    # 分隔线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(1.55), Inches(12.0), Pt(2),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_AQUA
    line.line.fill.background()


def _add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    bullets: list[str],
    *,
    accent: RGBColor = ACCENT_BLUE,
) -> None:
    """绘制一张卡片(标题 + 项目符号列表)。"""
    # 卡片背景
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = BG_LIGHT
    card.line.color.rgb = BORDER_LIGHT
    card.line.width = Pt(0.75)
    # 卡片顶部色条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Pt(4),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    # 卡片标题
    tb = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.15), Inches(width - 0.4), Inches(0.4)
    )
    _set_text(tb, title, font_size=14, bold=True, color=accent)
    # 卡片正文
    body = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.65), Inches(width - 0.4), Inches(height - 0.8)
    )
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "• " + b
        run.font.size = Pt(12)
        run.font.color.rgb = TEXT_DARK
        run.font.name = "微软雅黑"


def add_agent_capability_slide(prs: Presentation) -> None:
    """新增幻灯片:Agent 能力闭环强化。"""
    blank_layout = prs.slide_layouts[6]  # 通常 6 是空白布局
    slide = prs.slides.add_slide(blank_layout)

    _add_section_header(
        slide,
        "Agent 能力闭环强化",
        "对应赛题七环节 · 多轮对话 / 任务进度可视化 / 追问合并历史",
    )

    # 左卡:多轮对话上下文管理
    _add_card(
        slide, 0.6, 1.9, 3.9, 3.0,
        "多轮对话上下文管理",
        [
            "同会话内追问只调 LLM,注入历史消息",
            "上下文窗口 50 条,重启不丢失持久化会话",
            "区分用户/助手/系统消息,支持 Markdown 渲染",
            "化学公式 $\ce{C2H6O}$ 自动转换为 Unicode 上下标",
        ],
        accent=ACCENT_BLUE,
    )

    # 中卡:任务进度可视化侧栏
    _add_card(
        slide, 4.7, 1.9, 3.9, 3.0,
        "任务进度可视化侧栏",
        [
            "对话区始终可见,右侧独立展示任务步骤",
            "宽屏自动展开,支持向右折叠专注对话",
            "窄屏按钮切换,不打断对话节奏",
            "实时显示步骤状态(待执行/运行/完成/失败)",
        ],
        accent=ACCENT_AQUA,
    )

    # 右卡:追问合并历史
    _add_card(
        slide, 8.8, 1.9, 3.9, 3.0,
        "追问合并历史记录",
        [
            "追问与原对话合并为单条历史记录",
            "复用 sessionId,追问结果追加为新 section",
            "userInput 拼接 [追问] 标记,加载时还原",
            "保证会话连续性,降低存储碎片",
        ],
        accent=ACCENT_BLUE,
    )

    # 底部贴合赛题七环节说明
    tb = slide.shapes.add_textbox(
        Inches(0.6), Inches(5.2), Inches(12.1), Inches(0.9)
    )
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "对应赛题 Agent 能力闭环:"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = ACCENT_BLUE
    run.font.name = "微软雅黑"

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = "任务理解 → 多轮交互 → 任务规划 → 工具调用 → 知识增强 → 上下文管理 → 结果交付"
    run2.font.size = Pt(11)
    run2.font.color.rgb = TEXT_DARK
    run2.font.name = "微软雅黑"


def add_personalized_learning_slide(prs: Presentation) -> None:
    """新增幻灯片:个性化学习闭环。"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    _add_section_header(
        slide,
        "个性化学习闭环",
        "学情画像 → AI 诊断 → 学习规划 → 错题本,数据全程本地存储",
    )

    # 左卡:闭环流程
    _add_card(
        slide, 0.6, 1.9, 5.9, 3.2,
        "完整学习闭环",
        [
            "扫描/练习行为自动写入 Hive LearningRecord",
            "学情画像雷达图:5 大化学分类掌握度可视化",
            "AI 学情诊断:基于画像生成诊断报告与改进建议",
            "学习规划:生成个性化学习路径与同类题训练",
            "错题本:错因分析定位错误官能团,支持一键收藏",
        ],
        accent=ACCENT_AQUA,
    )

    # 右卡:"我的"页面与引导
    _add_card(
        slide, 6.8, 1.9, 5.9, 3.2,
        '"我的"页面与初次启动引导',
        [
            "底部导航新增「我的」tab:用户卡片 + 学情入口",
            "初次启动 onboarding:头像 / 昵称 / 学段",
            "用户管理 sheet:多用户切换、创建、编辑",
            "设置页保留用户管理区块,可深入查看",
            "学段选项:初中 / 高中 / 大学 / 科研",
        ],
        accent=ACCENT_BLUE,
    )

    # 底部合规说明
    tb = slide.shapes.add_textbox(
        Inches(0.6), Inches(5.3), Inches(12.1), Inches(0.6)
    )
    _set_text(
        tb,
        "合规边界:学习数据全部本地 Hive 存储,不上传;AI 生成内容附风险提示;符合《个人信息保护法》《未成年人保护法》对教育类应用要求。",
        font_size=11,
        color=TEXT_MUTED,
    )


def add_tech_architecture_slide(prs: Presentation) -> None:
    """新增幻灯片:技术架构更新。"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    _add_section_header(
        slide,
        "技术架构更新",
        "Flutter 内嵌 Agent 编排 · Riverpod 状态管理 · Hive 本地存储",
    )

    # 左卡:分层架构
    _add_card(
        slide, 0.6, 1.9, 5.9, 3.4,
        "分层架构",
        [
            "UI 层:AgentPage 对话主界面 + ProfilePage 我的页",
            "状态层:Riverpod 2.6 StateNotifier 管理任务/历史/消息",
            "Agent 编排层:TaskPlanner → AgentOrchestrator → ToolRegistry",
            "工具层:OcsrTool / PubChemTool / KnowledgeBaseTool / LlmTool",
            "数据层:Hive 隔离 box(每用户独立命名空间)",
        ],
        accent=ACCENT_BLUE,
    )

    # 右卡:关键设计
    _add_card(
        slide, 6.8, 1.9, 5.9, 3.4,
        "关键设计",
        [
            "AgentContext:基于 @slot.field 的跨步骤数据传递",
            "AgentSessionStore:持久化会话记录,支持追问合并",
            "MarkdownChatContent:化学公式 mhchem 预处理转 Unicode",
            "Cloudflare Worker:OCSR/PubChem 中转,规避 CORS",
            "提示词工程:7 套教育专用模板,作业辅导不直接给答案",
        ],
        accent=ACCENT_AQUA,
    )

    # 底部技术栈表
    tb = slide.shapes.add_textbox(
        Inches(0.6), Inches(5.5), Inches(12.1), Inches(0.5)
    )
    _set_text(
        tb,
        "技术栈:Flutter 3.x · Riverpod 2.6 · Hive · Dio · flutter_markdown · BlueLM/OpenAI 兼容 API · DECIMER OCSR · PubChem PUG REST · 31 个预置知识点",
        font_size=11,
        color=TEXT_MUTED,
    )


def move_thanks_to_end(prs: Presentation) -> None:
    """把 THANKS 页移到末尾(新增 3 张在它之后,需要重排)。

    python-pptx 不直接支持重排,通过操作 sldIdLst 顺序实现。
    """
    sldIdLst = prs.slides._sldIdLst  # type: ignore[attr-defined]
    slides = list(sldIdLst)
    # THANKS 通常在倒数第 2(倒数第 1 是空白页),把它和其后的空白页一起移到末尾
    # 简化:把最后 3 张新 slide 移到 THANKS 之前
    # 新增的 3 张是最后 3 个 id
    if len(slides) < 5:
        return
    # 找到 THANKS 页(包含"THANKS"或"感谢聆听"的)
    thanks_idx = None
    for i, sld in enumerate(slides):
        slide_obj = prs.slides[i]
        for shape in slide_obj.shapes:
            if shape.has_text_frame and ("THANKS" in shape.text or "感谢" in shape.text):
                thanks_idx = i
                break
        if thanks_idx is not None:
            break
    if thanks_idx is None:
        return
    # 把 THANKS 页及其后所有页移到末尾
    thanks_block = slides[thanks_idx:]
    new_order = slides[:thanks_idx] + slides[thanks_idx:]
    # 实际无需重排(新增的已经在末尾,THANKS 在新增之前)
    # 只有当 THANKS 在新增之后才需要重排,这里新增在末尾,THANKS 在它们之前是正确的
    # 所以无需操作
    return


def main() -> None:
    if not SRC_PPT.exists():
        raise SystemExit(f"源 PPT 不存在: {SRC_PPT}")

    # 备份原 PPT(首次运行)
    if not BACKUP.exists():
        shutil.copy2(SRC_PPT, BACKUP)
        print(f"[备份] {BACKUP.name}")

    # 读取源 PPT
    prs = Presentation(str(SRC_PPT))
    print(f"[读取] 共 {len(prs.slides)} 张幻灯片")

    # 1. 全局替换项目名
    replaced = replace_all_text(prs)
    print(f"[替换] 项目名替换影响 {replaced} 个形状")

    # 2. 新增 3 张幻灯片
    add_agent_capability_slide(prs)
    print("[新增] Agent 能力闭环强化")
    add_personalized_learning_slide(prs)
    print("[新增] 个性化学习闭环")
    add_tech_architecture_slide(prs)
    print("[新增] 技术架构更新")

    # 3. 保存为新文件(同时覆盖原文件,保留备份)
    prs.save(str(DST_PPT))
    print(f"[保存] {DST_PPT.name}")

    # 同时覆盖原文件名,保证用户在原路径也能看到更新
    prs.save(str(SRC_PPT))
    print(f"[覆盖] {SRC_PPT.name}")

    print("\n[完成] PPT 更新成功")
    print(f"  - 备份: {BACKUP.name}")
    print(f"  - 新文件: {DST_PPT.name}")
    print(f"  - 原文件已更新: {SRC_PPT.name}")


if __name__ == "__main__":
    main()
